"""Render an ArtifactSpec to a PowerPoint (.pptx) deck: a title slide + one slide per section."""

from typing import Optional

from pptx import Presentation

from src.app.core.artifacts.spec import ArtifactSpec, Template, claim_suffix


def render_pptx(spec: ArtifactSpec, path: str, template: Optional[Template] = None) -> str:
    """Write ``spec`` as a .pptx deck at ``path`` and return the path.

    A title slide, then one bulleted slide per section; each bullet carries its traceability
    suffix (source or unsourced marker).
    """
    template = template or Template()
    prs = Presentation()

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = spec.title
    if spec.subtitle and len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = spec.subtitle

    for section in spec.sections:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = section.heading
        body = slide.placeholders[1].text_frame
        body.clear()
        for i, claim in enumerate(section.claims):
            line = f"{claim.text}{claim_suffix(claim)}"
            paragraph = body.paragraphs[0] if i == 0 else body.add_paragraph()
            paragraph.text = line

    prs.save(path)
    return path
