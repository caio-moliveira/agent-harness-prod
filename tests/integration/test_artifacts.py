"""Integration tests for artifact generation + traceability (#18, RF-11/12/13).

Renders real .docx/.pptx files and re-opens them to assert content, that each claim shows its
source, and that an unsourced claim is flagged (not dropped). Also checks the template decouples
visual identity and that generation records an audit event.
"""

import pytest
from httpx import AsyncClient

from src.app.core.provenance import Source

pytestmark = pytest.mark.asyncio


def _spec():
    from src.app.core.artifacts import ArtifactSpec, Claim, Section

    return ArtifactSpec(
        title="Análise de Variação Orçamentária",
        subtitle="Fechamento mensal",
        sections=[
            Section(
                heading="Resultados",
                claims=[
                    Claim(
                        text="Receita realizada superou o planejado em 8%.",
                        source=Source.from_query(sql="SELECT SUM(receita) FROM realizado", tables=["realizado"]),
                    ),
                    Claim(text="A tendência sugere aceleração no próximo trimestre."),  # no source
                ],
            )
        ],
    )


class TestDocxRender:
    def test_docx_has_content_and_traceability(self, tmp_path):
        from docx import Document

        from src.app.core.artifacts import render_docx

        path = str(tmp_path / "relatorio.docx")
        render_docx(_spec(), path)

        text = "\n".join(p.text for p in Document(path).paragraphs)
        assert "Análise de Variação Orçamentária" in text
        assert "Resultados" in text
        assert "superou o planejado" in text
        assert "[Fonte:" in text and "realizado" in text  # sourced claim cites its source
        assert "[SEM FONTE]" in text                        # unsourced claim is flagged, not dropped

    def test_template_changes_only_style_not_content(self, tmp_path):
        from docx import Document

        from src.app.core.artifacts import Template, render_docx

        path = str(tmp_path / "branded.docx")
        render_docx(_spec(), path, template=Template(name="acme", primary_color="C00000"))
        text = "\n".join(p.text for p in Document(path).paragraphs)
        assert "superou o planejado" in text  # same content under a different template


class TestPptxRender:
    def test_pptx_has_title_and_section_slides(self, tmp_path):
        from pptx import Presentation

        from src.app.core.artifacts import render_pptx

        path = str(tmp_path / "deck.pptx")
        render_pptx(_spec(), path)

        prs = Presentation(path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text)
        joined = "\n".join(texts)
        assert "Análise de Variação Orçamentária" in joined
        assert "Resultados" in joined
        assert "[Fonte:" in joined
        assert "[SEM FONTE]" in joined


class TestTraceability:
    def test_unsourced_claims_detected(self):
        from src.app.core.artifacts import unsourced_claims

        missing = unsourced_claims(_spec())
        assert len(missing) == 1
        assert "aceleração" in missing[0].text


class TestGenerateAndAudit:
    async def test_generate_records_artifact_event(self, client: AsyncClient, tmp_path):
        from src.app.core.artifacts import generate_artifact
        from src.app.core.session.event_repository import SessionEventRepository

        path = str(tmp_path / "out.docx")
        await generate_artifact(_spec(), "docx", path, user_id=1, agent_id=7, session_id="sess-art")

        events = await SessionEventRepository().get_session_events("sess-art")
        assert len(events) == 1
        assert events[0].event_type == "artifact_generated"
        assert events[0].payload["format"] == "docx"
        assert events[0].payload["unsourced"] == 1

    async def test_unsupported_format_raises(self, tmp_path):
        from src.app.core.artifacts import generate_artifact

        with pytest.raises(ValueError):
            await generate_artifact(_spec(), "pdf", str(tmp_path / "x.pdf"))
